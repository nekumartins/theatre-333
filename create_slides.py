from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
PURPLE = RgbColor(139, 92, 246)
PINK = RgbColor(236, 72, 153)
DARK_BG = RgbColor(10, 10, 15)
CARD_BG = RgbColor(18, 18, 26)
WHITE = RgbColor(255, 255, 255)
GRAY = RgbColor(156, 163, 175)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    
    # Content
    y_pos = 1.5
    for item in content_items:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.333), Inches(0.5))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        y_pos += 0.6
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items, left_title="", right_title=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    
    # Left column title
    if left_title:
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(0.5))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PINK
    
    # Left content
    y_pos = 2.2 if left_title else 1.5
    for item in left_items:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(6), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        y_pos += 0.5
    
    # Right column title
    if right_title:
        right_title_box = slide.shapes.add_textbox(Inches(6.833), Inches(1.5), Inches(6), Inches(0.5))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PINK
    
    # Right content
    y_pos = 2.2 if right_title else 1.5
    for item in right_items:
        box = slide.shapes.add_textbox(Inches(6.833), Inches(y_pos), Inches(6), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        y_pos += 0.5
    
    return slide

# ============= CREATE SLIDES =============

# Slide 1: Title
slide = add_title_slide(prs, "Theatre Booking System", "A Complete Database Design & Implementation Project")
# Add additional info
info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(2))
tf = info_box.text_frame
p = tf.paragraphs[0]
p.text = "CSC 333 - Database Systems"
p.font.size = Pt(20)
p.font.color.rgb = PINK
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nPresented by: Chukwuneku Akpotohwo"
p.font.size = Pt(18)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "November 2025"
p.font.size = Pt(16)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

# Slide 2: Project Overview
add_two_column_slide(prs, "Project Overview",
    ["Design comprehensive database for theatre operations",
     "Online ticket booking platform with real-time features",
     "User registration, authentication, and authorization",
     "Complete booking workflow from selection to e-ticket"],
    ["Backend: FastAPI (Python)",
     "Database: MySQL 8.0",
     "ORM: SQLAlchemy with Alembic migrations",
     "Frontend: HTML, Tailwind CSS, JavaScript"],
    "Objectives", "Technology Stack")

# Slide 3: Business Requirements
slide = add_content_slide(prs, "Business Requirements", [])
requirements = [
    ("👥 User Management", "Registration, login, profiles, role-based access"),
    ("🎭 Show Management", "Multiple shows, genres, venues, scheduling"),
    ("💺 Seat Selection", "Interactive seat map, real-time availability, pricing tiers"),
    ("🎫 Booking System", "Multi-seat booking, confirmation, cancellation"),
    ("💳 Payment Processing", "Multiple methods, secure transactions, refunds"),
    ("🎟️ E-Ticket System", "QR codes, digital delivery, validation")
]
y = 1.5
for emoji_title, desc in requirements:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.333), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = emoji_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    y += 0.9

# Slide 4: Conceptual Model (ERD)
slide = add_content_slide(prs, "Conceptual Model (ERD)", [])
entities = [
    "USER (user_id PK, email, password_hash, full_name, phone, is_admin)",
    "SHOW (show_id PK, title, description, duration, genre_id FK, venue_id FK)",
    "PERFORMANCE (performance_id PK, show_id FK, performance_date, base_price)",
    "BOOKING (booking_id PK, user_id FK, performance_id FK, total_amount, status)",
    "PAYMENT (payment_id PK, booking_id FK, amount, method, status, transaction_ref)",
    "SEAT (seat_id PK, row_name, seat_number, category_id FK)",
    "BOOKING_SEAT (booking_seat_id PK, booking_id FK, seat_id FK, price_paid)"
]
y = 1.8
for entity in entities:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.333), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = entity
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    y += 0.6

# Add legend
legend_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
tf = legend_box.text_frame
p = tf.paragraphs[0]
p.text = "PK = Primary Key    FK = Foreign Key"
p.font.size = Pt(14)
p.font.color.rgb = GRAY

# Slide 5: Entity Relationships
slide = add_content_slide(prs, "Entity Relationships", [])
relationships = [
    ("USER → BOOKING", "1:N", "A user can make many bookings"),
    ("SHOW → PERFORMANCE", "1:N", "A show has many performances"),
    ("PERFORMANCE → BOOKING", "1:N", "A performance has many bookings"),
    ("BOOKING → BOOKING_SEAT", "1:N", "A booking includes many seats"),
    ("BOOKING → PAYMENT", "1:1", "Each booking has one payment"),
    ("VENUE → SHOW", "1:N", "A venue hosts many shows"),
    ("SEAT_CATEGORY → SEAT", "1:N", "A category contains many seats")
]
y = 1.8
for rel, cardinality, desc in relationships:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(5), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = rel
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    
    box2 = slide.shapes.add_textbox(Inches(5.5), Inches(y), Inches(1.5), Inches(0.5))
    tf = box2.text_frame
    p = tf.paragraphs[0]
    p.text = cardinality
    p.font.size = Pt(18)
    p.font.color.rgb = PINK
    
    box3 = slide.shapes.add_textbox(Inches(7), Inches(y), Inches(6), Inches(0.5))
    tf = box3.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    y += 0.7

# Slide 6: Normalization
slide = add_content_slide(prs, "Normalization Process", [])
nf_items = [
    ("✓ First Normal Form (1NF)", [
        "All attributes contain atomic values",
        "Each table has a primary key",
        "No repeating groups (seats in separate table)"
    ]),
    ("✓ Second Normal Form (2NF)", [
        "Satisfies 1NF requirements",
        "No partial dependencies on composite keys",
        "All non-key attributes fully depend on primary key"
    ]),
    ("✓ Third Normal Form (3NF)", [
        "Satisfies 2NF requirements", 
        "No transitive dependencies",
        "Genre, Venue, SeatCategory extracted to separate tables"
    ])
]
y = 1.5
for title, items in nf_items:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.333), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RgbColor(34, 197, 94)  # Green
    y += 0.5
    for item in items:
        box2 = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.4))
        tf = box2.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        y += 0.4
    y += 0.3

# Slide 7: Physical Model (SQL)
slide = add_content_slide(prs, "Physical Model (MySQL)", [])
sql_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5))
tf = sql_box.text_frame
code = """CREATE TABLE users (
    user_id INT PRIMARY KEY AUTO_INCREMENT,
    email VARCHAR(255) UNIQUE NOT NULL,
    password_hash VARCHAR(255) NOT NULL,
    full_name VARCHAR(100),
    is_admin BOOLEAN DEFAULT FALSE
);

CREATE TABLE bookings (
    booking_id INT PRIMARY KEY AUTO_INCREMENT,
    user_id INT NOT NULL,
    performance_id INT NOT NULL,
    total_amount DECIMAL(10,2) NOT NULL,
    status ENUM('pending','confirmed','cancelled'),
    FOREIGN KEY (user_id) REFERENCES users(user_id),
    FOREIGN KEY (performance_id) REFERENCES performances(performance_id)
);"""
p = tf.paragraphs[0]
p.text = code
p.font.size = Pt(12)
p.font.name = "Consolas"
p.font.color.rgb = WHITE

# Stats
stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.5))
tf = stats_box.text_frame
p = tf.paragraphs[0]
p.text = "11 Tables  |  MySQL 8.0  |  15+ Foreign Keys  |  InnoDB Engine"
p.font.size = Pt(16)
p.font.color.rgb = PINK
p.alignment = PP_ALIGN.CENTER

# Slide 8: Application Architecture
slide = add_content_slide(prs, "Application Architecture", [])
layers = [
    ("🖥️ Frontend", "HTML/Jinja2, Tailwind CSS, JavaScript"),
    ("⚡ API Layer", "FastAPI, REST Endpoints, JWT Authentication"),
    ("🔄 ORM Layer", "SQLAlchemy, Alembic Migrations"),
    ("🗄️ Database", "MySQL 8.0, InnoDB, ACID Compliant")
]
y = 2
for title, desc in layers:
    box = slide.shapes.add_textbox(Inches(2), Inches(y), Inches(9), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    y += 1.1

# Slide 9: Key Features
add_two_column_slide(prs, "Key Features",
    ["Secure password hashing (bcrypt)",
     "JWT token-based authentication",
     "Token expiration & validation",
     "Protected API routes",
     "Role-based access control (RBAC)"],
    ["Interactive seat map (491 seats)",
     "5 seat categories with pricing",
     "Real-time availability updates",
     "Multi-seat selection",
     "QR code e-tickets"],
    "🔐 Authentication", "💺 Seat Selection")

# Slide 10: Sample Queries
slide = add_content_slide(prs, "Sample SQL Queries", [])
query1 = """-- Get User's Booking History
SELECT b.booking_id, s.title, p.performance_date, b.total_amount, b.status
FROM bookings b
JOIN performances p ON b.performance_id = p.performance_id
JOIN shows s ON p.show_id = s.show_id
WHERE b.user_id = ? ORDER BY b.booking_date DESC;"""

query2 = """-- Get Available Seats for Performance
SELECT s.seat_id, s.row_name, s.seat_number, sc.category_name
FROM seats s JOIN seat_categories sc ON s.category_id = sc.category_id
WHERE s.seat_id NOT IN (
    SELECT bs.seat_id FROM booking_seats bs
    JOIN bookings b ON bs.booking_id = b.booking_id
    WHERE b.performance_id = ? AND b.status != 'cancelled'
);"""

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(2.5))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = query1
p.font.size = Pt(11)
p.font.name = "Consolas"
p.font.color.rgb = WHITE

box2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(2.5))
tf = box2.text_frame
p = tf.paragraphs[0]
p.text = query2
p.font.size = Pt(11)
p.font.name = "Consolas"
p.font.color.rgb = WHITE

# Slide 11: Security
add_two_column_slide(prs, "Security Implementation",
    ["Bcrypt hashing algorithm",
     "Salt rounds for extra security",
     "No plain-text password storage",
     "JWT token-based sessions",
     "Token expiration management"],
    ["Admin vs User roles",
     "Protected admin routes",
     "SQL injection prevention (ORM)",
     "Input validation (Pydantic)",
     "CORS configuration"],
    "🔒 Password & Auth", "🛡️ Access Control")

# Slide 12: Testing
slide = add_content_slide(prs, "Testing & Validation", [])
test_results = [
    "✓ All Frontend Pages: 200 OK (Home, Shows, Login, Register, Bookings, Admin)",
    "✓ All API Endpoints: 200 OK (/api/shows, /api/performances, /api/bookings, /api/payments)",
    "✓ Authentication Flow: Login, Token Generation, Protected Routes",
    "✓ Booking Flow: Seat Selection → Payment → Confirmation → E-Ticket",
    "✓ Database Integrity: Foreign Key Constraints, Transactions"
]
y = 1.8
for item in test_results:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.333), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(34, 197, 94)  # Green
    y += 0.7

# Stats
stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1))
tf = stats_box.text_frame
p = tf.paragraphs[0]
p.text = "100% Endpoint Coverage  |  0 Critical Errors  |  Production Ready"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = PINK
p.alignment = PP_ALIGN.CENTER

# Slide 13: Challenges & Solutions
slide = add_content_slide(prs, "Challenges & Solutions", [])
challenges = [
    ("Real-time Seat Availability", "Database transactions with proper locking"),
    ("Token Authentication Across Pages", "Centralized localStorage key management"),
    ("Complex Seat Pricing", "Price multiplier system with seat categories"),
    ("Performance Route Missing", "Added dedicated performance detail endpoint"),
    ("Database Migration", "SQLite to MySQL with Alembic migrations")
]
y = 1.8
for challenge, solution in challenges:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(5.5), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "❌ " + challenge
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(248, 113, 113)  # Red
    
    box2 = slide.shapes.add_textbox(Inches(6.5), Inches(y), Inches(6), Inches(0.5))
    tf = box2.text_frame
    p = tf.paragraphs[0]
    p.text = "✓ " + solution
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(34, 197, 94)  # Green
    y += 0.8

# Slide 14: Conclusion
slide = add_title_slide(prs, "Conclusion", "")
summary_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
tf = summary_box.text_frame
p = tf.paragraphs[0]
p.text = "Successfully designed and implemented a comprehensive theatre booking system demonstrating proper database design principles, normalization, and full-stack development practices."
p.font.size = Pt(20)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Stats
stats = ["11 Database Tables", "3NF Normalization", "15+ API Endpoints", "100% Functional"]
x = 1.5
y = 4
for stat in stats:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.5), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = stat
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER
    x += 2.8

# Slide 15: Thank You
slide = add_title_slide(prs, "Thank You!", "Questions?")
github_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1))
tf = github_box.text_frame
p = tf.paragraphs[0]
p.text = "GitHub: github.com/nekumartins/theatre-333"
p.font.size = Pt(16)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "Email: akpotohwoo@gmail.com"
p.font.size = Pt(16)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

# Save
prs.save('presentation/Theatre_Booking_System.pptx')
print("✓ Presentation saved to: presentation/Theatre_Booking_System.pptx")
print("\nTo use in Google Slides:")
print("1. Go to slides.google.com")
print("2. Click 'Blank' to create new presentation")
print("3. File → Import slides → Upload → Select the .pptx file")
print("4. Select all slides and click 'Import slides'")
